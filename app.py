import os
import fitz
import docx
from pptx import Presentation
from datetime import datetime
import streamlit as st
import streamlit.components.v1 as components
from dotenv import load_dotenv
from openai import OpenAI
from supabase_client import (
    supabase,
    sign_in,
    sign_up,
    get_user,
    sign_out,
    sign_in_with_google,
)

# ---------------- SETUP ----------------
load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
st.set_page_config(page_title="AI Study Hub", page_icon="🧠", layout="wide")

# Utility to rerun safely across Streamlit versions
def rerun():
    if hasattr(st, "rerun"):
        st.rerun()
    elif hasattr(st, "experimental_rerun"):
        st.experimental_rerun()

# ---------------- AUTH STATE ----------------
if "user" not in st.session_state:
    st.session_state.user = None

# ---------------- LOGIN UI ----------------
def login_ui():
    st.title("🔐 Welcome to AI Study Hub")
    tab1, tab2 = st.tabs(["Login", "Sign Up"])

    with tab1:
        email = st.text_input("Email", key="login_email")
        password = st.text_input("Password", type="password", key="login_password")
        if st.button("Login", key="login_btn"):
            try:
                sign_in(email, password)
                user = get_user()
                if user:
                    st.session_state.user = user
                    st.success(f"Welcome back, {email}!")
                    rerun()
                else:
                    st.error("Invalid credentials or user not found.")
            except Exception as e:
                st.error(f"Login failed: {e}")

    with tab2:
        email = st.text_input("Email", key="signup_email")
        password = st.text_input("Password", type="password", key="signup_password")
        if st.button("Create Account", key="signup_btn"):
            try:
                res = sign_up(email, password)
                if res.user:
                    st.success("✅ Account created! Check your email for verification.")
                else:
                    st.warning("Sign-up may have failed; try again.")
            except Exception as e:
                st.error(f"Sign-up failed: {e}")

    st.markdown("### 🔓 Or Sign in with Google")
    if st.button("Continue with Google", type="primary"):
        try:
            res = sign_in_with_google()
            st.success("Redirecting to Google sign-in...")
            st.markdown(f"[👉 Click here if not redirected automatically]({res.url})")
        except Exception as e:
            st.error(f"Google login failed: {e}")

# ---------------- LOGIN CHECK ----------------
user = get_user()
if not st.session_state.user and not user:
    login_ui()
    st.stop()
else:
    if user:
        st.session_state.user = user

# ---------------- SESSION STATE INIT ----------------
defaults = {
    "show_library": False,
    "resumed_session": None,
    "quiz_data": None,
    "dark_mode": False,
}
for key, val in defaults.items():
    if key not in st.session_state:
        st.session_state[key] = val

# ---------------- SIDEBAR ----------------
st.sidebar.header("Account")
if st.session_state.user:
    st.sidebar.success(f"Logged in as: {st.session_state.user.email}")
    if st.sidebar.button("🚪 Logout"):
        sign_out()
        st.session_state.user = None
        rerun()
else:
    st.sidebar.info("Not logged in")

# ---------------- HEADER ----------------
st.title("🧠 AI Flashcard Generator")
st.caption("Upload your study materials, generate flashcards, and master your topics.")
st.markdown("### 📘 Start a New Study Session")
st.write("Upload a file or paste text below, then click **Generate Flashcards** to begin!")

# ---------------- FILE / TEXT INPUT ----------------
uploaded_file = st.file_uploader(
    "Upload a file (PDF, DOCX, PPTX, TXT)", type=["pdf", "docx", "pptx", "txt"]
)
user_text = st.text_area("Or paste text here:", height=200)
input_text = ""

def extract_text(file):
    ext = file.name.split(".")[-1].lower()
    if ext == "pdf":
        text = ""
        with fitz.open(stream=file.read(), filetype="pdf") as pdf:
            for p in pdf:
                text += p.get_text()
        return text
    elif ext == "docx":
        d = docx.Document(file)
        return "\n".join([p.text for p in d.paragraphs])
    elif ext == "pptx":
        prs = Presentation(file)
        return "\n".join(
            [sh.text for s in prs.slides for sh in s.shapes if hasattr(sh, "text")]
        )
    elif ext == "txt":
        return file.read().decode("utf-8")
    return ""

if uploaded_file:
    input_text = extract_text(uploaded_file)
    st.success(f"✅ Extracted text from **{uploaded_file.name}**")
elif user_text.strip():
    input_text = user_text.strip()
# ---------------- FLASHCARD GENERATION + SAVE ----------------

def parse_flashcards(raw: str):
    """Return list of {'q':..., 'a':...} from LLM output with Q:/A:/--- format."""
    cards = []
    if not raw:
        return cards
    for chunk in raw.split("---"):
        chunk = chunk.strip()
        if not chunk:
            continue
        # tolerate variations like 'Q )', 'Q.' etc.
        q_part, a_part = None, None
        if "A:" in chunk:
            q_part, a_part = chunk.split("A:", 1)
        elif "Answer:" in chunk:
            q_part, a_part = chunk.split("Answer:", 1)
        else:
            continue
        q_txt = q_part.replace("Q:", "").replace("Question:", "").strip()
        a_txt = a_part.strip()
        if q_txt and a_txt:
            cards.append({"q": q_txt, "a": a_txt})
    return cards

def render_card(q, a, cid, color="#2563eb"):
    components.html(
        f"""
        <div style="perspective:1000px;display:inline-block;margin:10px;">
          <div id="card-{cid}" onclick="this.classList.toggle('flipped')"
               style="width:280px;height:180px;text-align:center;
                      transition:transform 0.6s;transform-style:preserve-3d;
                      box-shadow:0 4px 8px rgba(0,0,0,0.2);
                      border-radius:12px;cursor:pointer;position:relative;">
            <div style="position:absolute;width:100%;height:100%;
                        backface-visibility:hidden;border-radius:12px;
                        background-color:#f9fafb;color:#111;
                        display:flex;justify-content:center;
                        align-items:center;padding:10px;overflow:auto;">
              <b>Q:</b>&nbsp;{q}
            </div>
            <div style="position:absolute;width:100%;height:100%;
                        backface-visibility:hidden;border-radius:12px;
                        background-color:{color};color:white;
                        transform:rotateY(180deg);
                        display:flex;justify-content:center;
                        align-items:center;padding:10px;overflow:auto;">
              <b>A:</b>&nbsp;{a}
            </div>
          </div>
          <script>
            const c=document.getElementById('card-{cid}');
            c.addEventListener('click',()=>{{
                c.style.transform=c.classList.contains('flipped')?
                  'rotateY(180deg)':'rotateY(0deg)';
            }});
          </script>
        </div>
        """,
        height=220,
    )

def insert_session_to_db(session_row: dict):
    """Insert session row into Supabase and return the inserted row (with id)."""
    res = supabase.table("sessions").insert(session_row).execute()
    # Supabase python client returns {'data': [...], 'count': None}
    if hasattr(res, "data") and res.data and isinstance(res.data, list):
        return res.data[0]
    # Some client versions return dict
    if isinstance(res, dict) and res.get("data"):
        return res["data"][0]
    return session_row

# The main button that generates flashcards and saves the session
if st.button("✨ Generate Flashcards", key="gen_flashcards"):
    if not input_text:
        st.warning("Please upload or enter text first.")
    else:
        with st.spinner("Generating flashcards..."):
            # Keep prompt tight + truncate to avoid hitting token limits on large docs
            prompt = f"""
            You are a helpful AI tutor. Create 8–12 concise Q&A flashcards from this study text.
            Use EXACTLY this format and nothing else:

            Q: [question]
            A: [answer]
            ---
            
            Study text (truncated):
            {input_text[:8000]}
            """
            try:
                response = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.5,
                )
                flashcards_text = response.choices[0].message.content
            except Exception as e:
                st.error(f"Flashcard generation failed: {e}")
                flashcards_text = ""

        cards = parse_flashcards(flashcards_text)
        if not cards:
            st.warning("I couldn't parse any flashcards from the response. Try a smaller input or different text.")
        else:
            st.success(f"✅ Generated {len(cards)} flashcards!")

            # Display flashcards
            for i, c in enumerate(cards):
                render_card(c["q"], c["a"], f"fc-{i}")

            # Save study session to Supabase and activate it
            try:
                user_id = st.session_state.user.id
                session_row = {
                    "user_id": user_id,
                    "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "file_name": uploaded_file.name if uploaded_file else "Pasted Text",
                    "flashcards": flashcards_text,
                    "quiz_data": None,
                    "score": None,
                    "adaptive_data": None,
                }
                inserted = insert_session_to_db(session_row)

                # Make the newly created session the active one for Quiz/Adaptive
                st.session_state.resumed_session = inserted
                st.info("💾 Study session saved and set as your active session. You can start a quiz now.")
            except Exception as e:
                st.warning(f"Session saved locally only (Supabase insert failed): {e}")
                st.session_state.resumed_session = session_row
# ---------------- STUDY LIBRARY (LIST + RESUME) ----------------

st.markdown("---")
st.header("📚 Study Library")

def fetch_sessions(user_id):
    """Return a list of session rows for the user, newest first."""
    try:
        # Order by id desc is reliable since id is identity
        res = supabase.table("sessions") \
            .select("*") \
            .eq("user_id", user_id) \
            .order("id", desc=True) \
            .execute()
        return res.data or []
    except Exception as e:
        st.error(f"Error loading sessions: {e}")
        return []

def set_latest_session_active_if_missing():
    """If no active session selected, auto-pick the latest (if exists)."""
    if st.session_state.resumed_session:
        return
    try:
        user_id = st.session_state.user.id
        latest = fetch_sessions(user_id)
        if latest:
            st.session_state.resumed_session = latest[0]
    except Exception:
        pass

# Auto-pick latest session on first load (helps Quiz/Adaptive pages)
set_latest_session_active_if_missing()

# Toggle to show/hide library list
col_l, col_r = st.columns([0.7, 0.3])
with col_l:
    st.write("Review your past study sessions and continue where you left off.")
with col_r:
    if st.button("📁 Open / Close Library", key="toggle_lib"):
        st.session_state.show_library = not st.session_state.show_library

if st.session_state.show_library:
    user_id = st.session_state.user.id
    sessions = fetch_sessions(user_id)

    if not sessions:
        st.info("You don’t have any saved study sessions yet. Generate flashcards above to create your first session.")
    else:
        for sess in sessions:
            with st.expander(f"📘 {sess.get('file_name','(no name)')} — {sess.get('timestamp','')}", expanded=False):
                st.write(f"**Score:** {sess.get('score', 'Not yet taken')}")
                # Preview first 3 cards as plain text (keeps UI light)
                preview = (sess.get("flashcards") or "").split("---")[:3]
                if preview:
                    st.caption("Preview:")
                    st.code("\n---\n".join([p.strip() for p in preview if p.strip()]), language="markdown")

                cols = st.columns([0.6, 0.4])
                with cols[0]:
                    if st.button(f"▶️ Resume", key=f"resume_{sess['id']}"):
                        st.session_state.resumed_session = sess
                        st.success(f"Resumed: {sess.get('file_name','(no name)')}")
                        rerun()
                with cols[1]:
                    # Optional: a quick “Set Active” alias without reloading
                    if st.button("Set Active (no reload)", key=f"setactive_{sess['id']}"):
                        st.session_state.resumed_session = sess
                        st.toast("Active session updated ✅", icon="✅")

# Show currently active session summary (always visible)
st.markdown("---")
st.subheader("🔁 Active Session")
active = st.session_state.resumed_session
if not active:
    st.info("No active session yet. Open the Library and click **Resume** on a session, or generate new flashcards.")
else:
    st.write(f"**File:** {active.get('file_name','(no name)')}  |  **When:** {active.get('timestamp','')}  |  **Score:** {active.get('score','N/A')}")
    # Show first card question as a quick reminder
    first_card = (active.get("flashcards") or "").split("---")[0].strip()
    if first_card:
        st.caption("First card preview:")
        st.code(first_card, language="markdown")
# ---------------- QUIZ MODE ----------------
st.markdown("---")
st.header("🧩 Quiz Mode")

def generate_quiz_from_flashcards(flashcards_text: str) -> str:
    """Ask the model to build MCQs from flashcards."""
    prompt = f"""
    Create 5 multiple-choice questions from these flashcards:
    {flashcards_text}

    Format exactly like this:
    Q: What is ...?
    A) option1
    B) option2
    C) option3
    D) option4
    Correct: [A/B/C/D]
    ---
    """
    with st.spinner("Generating quiz questions..."):
        try:
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.6,
            )
            return resp.choices[0].message.content
        except Exception as e:
            st.error(f"Quiz generation failed: {e}")
            return ""

def parse_quiz(raw_text: str):
    """Parse Q/A/Correct blocks into list of dicts."""
    questions = []
    for chunk in raw_text.split("---"):
        if not chunk.strip():
            continue
        lines = [ln.strip() for ln in chunk.splitlines() if ln.strip()]
        q_text, options, correct = None, [], None
        for ln in lines:
            if ln.startswith("Q:"):
                q_text = ln[2:].strip()
            elif ln[:2] in ["A)", "B)", "C)", "D)"]:
                options.append(ln)
            elif ln.lower().startswith("correct"):
                correct = ln.split(":")[-1].strip().replace(".", "")
        if q_text and options and correct:
            questions.append({"question": q_text, "options": options, "correct": correct})
    return questions

# ---------------- QUIZ STATE ----------------
if "quiz_questions" not in st.session_state:
    st.session_state.quiz_questions = []
if "quiz_answers" not in st.session_state:
    st.session_state.quiz_answers = {}

active_sess = st.session_state.resumed_session
if not active_sess:
    st.info("ℹ️  You must have an active study session. Generate or resume one first.")
else:
    # Generate quiz
    if st.button("🧠 Start AI Quiz", key="start_quiz"):
        text = active_sess.get("flashcards", "")
        if not text.strip():
            st.warning("This session has no flashcards yet.")
        else:
            raw_quiz = generate_quiz_from_flashcards(text)
            qs = parse_quiz(raw_quiz)
            if qs:
                st.session_state.quiz_questions = qs
                st.session_state.quiz_answers = {}
                st.success(f"✅  Generated {len(qs)} questions!")
                rerun()
            else:
                st.error("Could not parse quiz from AI output. Try again.")

# ---------------- QUIZ PLAY ----------------
qs = st.session_state.quiz_questions
if qs:
    st.subheader("📖 Quiz Time!")
    for i, q in enumerate(qs):
        st.markdown(f"**Q{i+1}: {q['question']}**")
        choice = st.radio(
            "Select an answer:",
            [opt for opt in q["options"]],
            key=f"quiz_choice_{i}",
            horizontal=True,
            index=None,
        )
        st.session_state.quiz_answers[i] = choice

    if st.button("📊 Submit Quiz", key="submit_quiz"):
        score, total = 0, len(qs)
        for i, q in enumerate(qs):
            ans = st.session_state.quiz_answers.get(i, "")
            correct_letter = q["correct"].strip().upper().replace(".", "")
            if ans.startswith(correct_letter):
                score += 1

        percent = round((score / total) * 100)
        st.success(f"🏁 You scored {percent}% ({score}/{total})")

        # Save to Supabase
        try:
            uid = st.session_state.user.id
            sid = active_sess.get("id")
            if sid:
                supabase.table("sessions").update(
                    {"score": str(percent)}
                ).eq("id", sid).eq("user_id", uid).execute()
                st.toast("Score saved ✅", icon="💾")
            else:
                st.toast("Session not yet synced to DB (local only)", icon="⚠️")
        except Exception as e:
            st.warning(f"Could not save quiz score: {e}")
# ---------------- ADAPTIVE FLASHCARDS (Flip Card UI) ----------------
st.markdown("---")
st.header("🎯 Adaptive Flashcards")

def generate_adaptive_flashcards(session_data):
    """Use GPT to create tougher follow-up flashcards."""
    flashcards_text = session_data.get("flashcards", "")
    if not flashcards_text:
        return ""
    prompt = f"""
    You are an adaptive learning AI. From the following flashcards, identify weak or confusing areas.
    Then generate 5–10 new flashcards focusing on those topics, phrased in a slightly more challenging way.

    {flashcards_text}

    Format:
    Q: [question]
    A: [answer]
    ---
    """
    with st.spinner("Creating adaptive flashcards..."):
        try:
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.7,
            )
            return resp.choices[0].message.content
        except Exception as e:
            st.error(f"Adaptive flashcard generation failed: {e}")
            return ""

def render_adaptive_card(q, a, cid, color="#16a34a"):
    components.html(
        f"""
        <div style="perspective:1000px;display:inline-block;margin:10px;">
          <div id="adaptive-{cid}" onclick="this.classList.toggle('flipped')"
               style="width:280px;height:180px;text-align:center;
                      transition:transform 0.6s;transform-style:preserve-3d;
                      box-shadow:0 4px 8px rgba(0,0,0,0.3);
                      border-radius:12px;cursor:pointer;position:relative;">
            <div style="position:absolute;width:100%;height:100%;
                        backface-visibility:hidden;border-radius:12px;
                        background-color:#f9fafb;color:#111;
                        display:flex;justify-content:center;
                        align-items:center;padding:10px;">
              <b>Q:</b>&nbsp;{q}
            </div>
            <div style="position:absolute;width:100%;height:100%;
                        backface-visibility:hidden;border-radius:12px;
                        background-color:{color};color:white;
                        transform:rotateY(180deg);
                        display:flex;justify-content:center;
                        align-items:center;padding:10px;">
              <b>A:</b>&nbsp;{a}
            </div>
          </div>
          <script>
            const c=document.getElementById('adaptive-{cid}');
            c.addEventListener('click',()=>{{
                c.style.transform=c.classList.contains('flipped')?
                  'rotateY(180deg)':'rotateY(0deg)';
            }});
          </script>
        </div>
        """,
        height=220,
    )

if st.button("🔁 Generate Adaptive Flashcards", key="adaptive_btn"):
    sess = st.session_state.resumed_session
    if not sess:
        st.warning("Select or generate a study session first.")
    else:
        adaptive_text = generate_adaptive_flashcards(sess)
        if adaptive_text:
            st.success("✅ Adaptive flashcards generated!")
            cards = adaptive_text.split("---")
            for i, c in enumerate(cards):
                if "A:" in c:
                    q, a = c.split("A:")
                    render_adaptive_card(q.replace("Q:", "").strip(), a.strip(), f"ad-{i}")
            try:
                supabase.table("sessions").update(
                    {"adaptive_data": adaptive_text}
                ).eq("id", sess.get("id")).execute()
                st.toast("Adaptive flashcards saved ✅", icon="💾")
            except Exception as e:
                st.warning(f"Could not save to Supabase: {e}")


# ---------------- SMART RECOMMENDATIONS ----------------
st.markdown("---")
st.header("🧭 Smart Recommendations")

def get_recommendations(user_id):
    """Analyze performance trends and suggest next study focus."""
    try:
        res = supabase.table("sessions").select("*").eq("user_id", user_id).execute()
        sessions = res.data or []
        if not sessions:
            st.info("You don’t have any completed sessions yet.")
            return

        study_summary = "\n".join(
            [f"{s['file_name']} — score {s.get('score', 'N/A')}" for s in sessions if s.get("score")]
        )

        if not study_summary.strip():
            st.info("Complete at least one quiz to receive recommendations.")
            return

        with st.spinner("Analyzing your progress..."):
            prompt = f"""
            You are a study coach analyzing the user's learning history.
            Review the quiz scores below and provide:
            - A summary of their weak points
            - 3 recommended topics or files to revisit
            - 2 encouragement tips to improve consistency

            Study History:
            {study_summary}

            Write as short bullet points with emojis.
            """
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.6,
            )
            st.subheader("📊 Personalized Study Plan")
            st.markdown(resp.choices[0].message.content)
    except Exception as e:
        st.error(f"Could not generate recommendations: {e}")

if st.button("🧠 Generate Smart Recommendations", key="smart_recs"):
    get_recommendations(st.session_state.user.id)

# ---------------- UI ENHANCEMENTS ----------------
st.markdown("---")
st.header("⚙️ Interface Options")

# Dark Mode Toggle
col1, col2 = st.columns(2)
with col1:
    dark = st.checkbox("🌙 Enable Dark Mode", value=st.session_state.dark_mode)
    if dark != st.session_state.dark_mode:
        st.session_state.dark_mode = dark
        theme = "dark" if dark else "light"
        st.toast(f"Switched to {theme} mode 🌗", icon="💡")
        rerun()

if st.session_state.dark_mode:
    st.markdown(
        """
        <style>
        body, .stApp { background-color: #111 !important; color: #eee !important; }
        .stButton>button { background-color: #333 !important; color: #fff !important; border: 1px solid #444; }
        .stTextArea textarea, .stTextInput input { background-color: #222 !important; color: #eee !important; }
        </style>
        """,
        unsafe_allow_html=True,
    )

# Sidebar Navigation
st.sidebar.markdown("---")
st.sidebar.subheader("📂 Navigation")

nav_pages = ["🏠 Home", "📚 Library", "🧩 Quiz", "🎯 Adaptive", "🧭 Recommendations"]
choice = st.sidebar.radio("Navigate to:", nav_pages, key="nav_choice")

if choice == "🏠 Home":
    st.info("Upload a file or paste text to generate flashcards.")
elif choice == "📚 Library":
    st.session_state.show_library = True
    rerun()
elif choice == "🧩 Quiz":
    st.info("Scroll to the quiz section below.")
elif choice == "🎯 Adaptive":
    st.info("Scroll to Adaptive Flashcards to practice tough topics.")
elif choice == "🧭 Recommendations":
    st.info("Scroll to Smart Recommendations for your personalized feedback.")

# Progress Tracker
st.markdown("---")
st.header("📈 Study Progress")

def calc_progress(user_id):
    try:
        data = supabase.table("sessions").select("*").eq("user_id", user_id).execute().data
        if not data:
            return 0
        completed = sum(1 for s in data if s.get("score"))
        return int((completed / len(data)) * 100)
    except Exception:
        return 0

if st.session_state.user:
    progress = calc_progress(st.session_state.user.id)
    st.progress(progress / 100)
    st.caption(f"🎓 You’ve completed {progress}% of your study sessions.")
else:
    st.info("Log in to track your progress.")
